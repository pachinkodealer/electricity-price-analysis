"""
Build a single-slide 48x36" academic POSTER as .pptx (python-pptx).
Mirrors poster/electricity_poster.html. Charts are placed as pictures (aspect
preserved automatically). Output: poster/electricity_poster.pptx
Run:  py build_poster_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent.parent
CH = ROOT / "charts"
OUT = Path(__file__).parent / "electricity_poster.pptx"
EMU_IN = 914400

INK, SOFT, FAINT = "161B22", "47515E", "7A8593"
ELEC, ELEC_D = "1F77B4", "12547F"
GAS, GAS_D = "C85E1B", "9C470F"
TINT, GAS_TINT, LINE, WHITE = "F1F5F9", "F7EAE0", "D8DEE6", "FFFFFF"
SERIF, SANS = "Cambria", "Calibri"

prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def rgb(h):
    return RGBColor.from_string(h)


def text(x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "sa" in para:
            p.space_after = Pt(para["sa"])
        if "sb" in para:
            p.space_before = Pt(para["sb"])
        if "ls" in para:
            p.line_spacing = para["ls"]
        for run in para["runs"]:
            r = p.add_run()
            r.text = run["t"]
            f = r.font
            f.name = run.get("font", SANS)
            f.size = Pt(run.get("sz", 16))
            f.bold = run.get("b", False)
            f.italic = run.get("i", False)
            f.color.rgb = rgb(run.get("c", INK))
    return tb


def rrect(x, y, w, h, fill, radius=0.06, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    if line:
        sh.line.color.rgb = rgb(line)
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def oval(x, y, d, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def pic(path, x, y, w):
    p = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return p.height / EMU_IN  # rendered height in inches


def R(t, **k):
    k["t"] = t
    return k


# ---------------- header ----------------
text(1.6, 1.4, 34, 0.5, [{"runs": [R("ENERGY ECONOMICS  ·  DATA INVESTIGATION", font=SANS, sz=15, b=True, c=ELEC_D)]}])
text(1.6, 1.9, 33.5, 2.5, [{"runs": [R("Why Your Electric Bill Stopped Following Fuel Prices",
                                       font=SERIF, sz=58, b=True, c=INK)], "ls": 1.0}])
text(1.6, 4.45, 33.5, 0.8, [{"runs": [R("US electricity prices, natural gas, and the datacenter demand era — 1978–2026",
                                        font=SERIF, sz=25, i=True, c=SOFT)]}])
text(37.0, 1.7, 9.4, 2.6, [
    {"runs": [R("Ian Lee", font=SANS, sz=17, b=True, c=INK)], "align": PP_ALIGN.RIGHT, "sa": 3},
    {"runs": [R("Independent data analysis", font=SANS, sz=15, c=SOFT)], "align": PP_ALIGN.RIGHT, "sa": 3},
    {"runs": [R("June 2026", font=SANS, sz=15, c=SOFT)], "align": PP_ALIGN.RIGHT, "sa": 3},
    {"runs": [R("Python · pandas · statsmodels", font=SANS, sz=13, c=FAINT)], "align": PP_ALIGN.RIGHT},
])

# thesis
text(1.6, 5.5, 44.8, 1.5, [{"runs": [
    R("For decades, US retail electricity tracked the price of natural gas. Around 2020 that link broke — ", font=SERIF, sz=23, c=INK),
    R("fuel got cheap again, but bills kept climbing", font=SERIF, sz=23, b=True, c=ELEC_D),
    R(" — and the timing lines up with the first sustained growth in US power demand in two decades: the AI datacenter buildout.", font=SERIF, sz=23, c=INK),
], "ls": 1.1}])

# ---------------- KPI strip ----------------
kpis = [
    ("+48%", "Retail electricity, since Jan 2020", "BLS", ELEC_D),
    ("−65%", "Natural gas, from its 2022 peak", "HENRY HUB", GAS_D),
    ("+19%", "Electricity since that gas peak — it rose while gas fell", "2022 → 2026", ELEC_D),
    ("2.6 pts/yr", "Inflation unexplained by fuel or CPI, since 2023", "RESIDUAL", ELEC_D),
]
ky, kh = 7.35, 2.85
kx0, kgap = 1.6, 0.5
kw = (44.8 - 3 * kgap) / 4
for i, (num, lab, unit, col) in enumerate(kpis):
    kx = kx0 + i * (kw + kgap)
    rrect(kx, ky, kw, kh, TINT, radius=0.05)
    nsz = 56 if len(num) <= 4 else 40
    text(kx + 0.45, ky + 0.28, kw - 0.9, 1.35, [{"runs": [R(num, font=SERIF, sz=nsz, b=True, c=col)]}], anchor=MSO_ANCHOR.MIDDLE)
    text(kx + 0.45, ky + 1.62, kw - 0.9, 0.8, [{"runs": [R(lab, font=SANS, sz=15, c=SOFT)], "ls": 1.0}])
    text(kx + 0.45, ky + 2.42, kw - 0.9, 0.35, [{"runs": [R(unit, font=SANS, sz=11, b=True, c=FAINT)]}])

# ---------------- columns ----------------
TOP = 10.5
c1x, c1w = 1.6, 12.0
c2x, c2w = 14.4, 16.3
c3x, c3w = 31.5, 14.9


def eyebrow(x, y, w, s):
    text(x, y, w, 0.4, [{"runs": [R(s, font=SANS, sz=14, b=True, c=ELEC_D)]}])


def fig(x, y, w, title, tag, img, cap_runs):
    text(x, y, w * 0.7, 0.45, [{"runs": [R(title, font=SERIF, sz=20, b=True, c=INK)]}])
    text(x + w * 0.55, y + 0.03, w * 0.45, 0.4, [{"runs": [R(tag, font=SANS, sz=12, b=True, c=FAINT)]}], anchor=MSO_ANCHOR.TOP)
    ih = pic(img, x, y + 0.55, w)
    cy = y + 0.55 + ih + 0.12
    text(x, cy, w, 1.2, [{"runs": cap_runs, "ls": 1.0}])
    return cy + 1.25


# --- column 1 ---
cy = TOP
eyebrow(c1x, cy, c1w, "MOTIVATION"); cy += 0.45
text(c1x, cy, c1w, 0.7, [{"runs": [R("The question", font=SERIF, sz=26, b=True, c=INK)]}]); cy += 0.75
text(c1x, cy, c1w, 1.3, [{"runs": [R("From 2014 to 2020, average US electricity prices were essentially flat. Since 2020 they have risen 48%.", font=SANS, sz=18, c=INK)], "ls": 1.05}]); cy += 1.65
text(c1x, cy, c1w, 1.8, [{"runs": [R("A jump like that usually has a simple cause: costlier fuel, or general inflation lifting everything. So — is this just fuel and inflation, or something genuinely new?", font=SANS, sz=15, c=SOFT)], "ls": 1.05}]); cy += 2.1

eyebrow(c1x, cy, c1w, "DATA & METHOD"); cy += 0.45
text(c1x, cy, c1w, 0.9, [{"runs": [R("47 years, three public series", font=SERIF, sz=22, b=True, c=INK)], "ls": 1.0}]); cy += 1.0
series = [
    ("Retail electricity — $/kWh", "APU000072610 · BLS · 1978–26"),
    ("Consumer Price Index", "CPIAUCSL · BLS · monthly"),
    ("Henry Hub gas — $/MMBtu", "MHHNGSP · EIA · 1997–26"),
]
sp = []
for nm, meta in series:
    sp.append({"runs": [R(nm, font=SANS, sz=14, b=True, c=INK)], "sa": 1})
    sp.append({"runs": [R(meta, font=SANS, sz=12, c=FAINT)], "sa": 7})
text(c1x, cy, c1w, 2.6, sp); cy += 2.7
steps = [
    "Benchmark electricity against overall inflation (CPI).",
    "Lag-correlate electricity vs. gas, year over year.",
    "Fit a two-factor model (gas + CPI); read what it can't explain.",
]
stp = []
for i, s in enumerate(steps):
    stp.append({"runs": [R(f"{i+1}   ", font=SERIF, sz=16, b=True, c=ELEC), R(s, font=SANS, sz=15, c=INK)], "sa": 8, "ls": 1.0})
text(c1x, cy, c1w, 2.4, stp); cy += 2.5
rrect(c1x, cy, c1w, 0.7, TINT, radius=0.18)
text(c1x, cy, c1w, 0.7, [{"runs": [R("Reproducible · live FRED endpoints · no API keys", font=SANS, sz=13, b=True, c=ELEC_D)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
cy += 0.7 + 0.85

eyebrow(c1x, cy, c1w, "CONCLUSIONS"); cy += 0.5
text(c1x, cy, c1w, 4.6, [
    {"runs": [R("The link between fuel and your bill is broken.", font=SANS, sz=15, b=True, c=INK),
              R(" In the shale decade, cheap gas meant flat prices; that cool-off is absent today.", font=SANS, sz=15, c=SOFT)], "sa": 9, "ls": 1.05},
    {"runs": [R("Whatever replaced fuel is growing, not fading.", font=SANS, sz=15, b=True, c=INK),
              R(" The honest claim isn't “unprecedented” — it's elevated, persistent, and attached to demand that keeps rising.", font=SANS, sz=15, c=SOFT)], "ls": 1.05},
]); cy += 4.6
cy1_end = cy

# --- column 2 ---
cy = TOP
cy = fig(c2x, cy, c2w, "The ratchet", "RESULT 01", CH / "gas_vs_electricity.png", [
    R("Both series indexed to Jan 2020 = 100. Gas spiked ~4× in 2022 and ", font=SANS, sz=15, c=SOFT),
    R("round-tripped to −65%", font=SANS, sz=15, b=True, c=INK),
    R("; retail electricity followed it up — and never came back down.", font=SANS, sz=15, c=SOFT),
])
cy += 0.25
cy2_end = fig(c2x, cy, c2w, "Not just inflation", "RESULT 02", CH / "electricity_money_chart.png", [
    R("Rolling 6-year % change vs. CPI. The shaded band is the 2010s ", font=SANS, sz=15, c=SOFT),
    R("shale decade", font=SANS, sz=15, b=True, c=INK),
    R(", when cheap gas kept prices flat. Today electricity is outrunning inflation — steeper than 96% of all 6-year periods since 1984.", font=SANS, sz=15, c=SOFT),
])

# --- column 3 ---
cy = TOP
cy = fig(c3x, cy, c3w, "The unexplained gap", "RESULT 03", CH / "residual_model.png", [
    R("A two-factor OLS model (6-month-lagged gas + CPI), trained only on pre-2023 data, predicts what electricity inflation ", font=SANS, sz=15, c=SOFT),
    R("should", font=SANS, sz=15, i=True, c=SOFT),
    R(" be. The residual sits near zero for most of history — then ", font=SANS, sz=15, c=SOFT),
    R("lifts to ~2.6 pts/yr after 2023", font=SANS, sz=15, b=True, c=INK),
    R(" and stays there.", font=SANS, sz=15, c=SOFT),
])
cy += 0.2
eyebrow(c3x, cy, c3w, "FINDINGS"); cy += 0.5
findings = [
    [R("Electricity ratchets.", font=SANS, sz=15, b=True, c=INK),
     R(" Follows gas up with a ~6-month lag (r = 0.53) but did not retrace when gas gave back its 2022 spike.", font=SANS, sz=15, c=SOFT)],
    [R("A persistent ~2.6 pts/yr", font=SANS, sz=15, b=True, c=INK),
     R(" of electricity inflation since 2023 is unexplained by fuel or CPI — tracking the datacenter-driven return of demand growth.", font=SANS, sz=15, c=SOFT)],
    [R("Elevated, not unprecedented.", font=SANS, sz=15, b=True, c=INK),
     R(" Prior residual episodes (2005–09) decayed; this one hasn't — and its driver is still growing.", font=SANS, sz=15, c=SOFT)],
]
for i, runs in enumerate(findings):
    oval(c3x, cy, 0.5, ELEC)
    text(c3x, cy - 0.02, 0.5, 0.5, [{"runs": [R(str(i + 1), font=SERIF, sz=18, b=True, c=WHITE)], "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(c3x + 0.72, cy - 0.05, c3w - 0.72, 1.4, [{"runs": runs, "ls": 1.02}])
    cy += 1.5
cy += 0.15
lh = 2.9
rrect(c3x, cy, c3w, lh, GAS_TINT, radius=0.05)
text(c3x + 0.35, cy + 0.28, c3w - 0.7, lh - 0.5, [
    {"runs": [R("STATED LIMITATIONS", font=SANS, sz=12, b=True, c=GAS_D)], "sa": 8},
    {"runs": [R("The residual is “not gas, not CPI” — it bundles datacenter demand with grid capex, plant retirements, and electrification. Not a clean causal estimate.", font=SANS, sz=13.5, c=SOFT)], "sa": 7, "ls": 1.0},
    {"runs": [R("Overlapping YoY series are heavily autocorrelated (Durbin–Watson ≈ 0.2): a descriptive decomposition, not formal inference.", font=SANS, sz=13.5, c=SOFT)], "ls": 1.0},
])
cy3_end = cy + lh

# ---------------- footer (positioned just below the tallest column) ----------------
fy = max(cy1_end, cy2_end, cy3_end) + 0.8
text(1.6, fy, 31, 2.3, [
    {"runs": [R("Not everyone uses AI.", font=SERIF, sz=33, b=True, c=INK)], "sa": 4, "ls": 1.0},
    {"runs": [R("But ", font=SERIF, sz=33, b=True, c=INK), R("everyone pays an electric bill.", font=SERIF, sz=33, b=True, c=GAS_D)], "ls": 1.0},
], anchor=MSO_ANCHOR.MIDDLE)
text(34.2, fy + 0.35, 12.2, 1.7, [
    {"runs": [R("Data: BLS & EIA via FRED. Reproducible end-to-end.", font=SANS, sz=13, c=SOFT)], "align": PP_ALIGN.RIGHT, "sa": 5},
    {"runs": [R("github.com/pachinkodealer/electricity-price-analysis", font=SANS, sz=13, b=True, c=ELEC_D)], "align": PP_ALIGN.RIGHT, "sa": 5},
    {"runs": [R("© Ian Lee 2026", font=SANS, sz=12, c=FAINT)], "align": PP_ALIGN.RIGHT},
])

prs.save(str(OUT))
print("Wrote", OUT, f"({OUT.stat().st_size/1024:.0f} KB)  slide 48x36 in")
